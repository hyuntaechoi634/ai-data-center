#!/usr/bin/env python3
from __future__ import annotations

import argparse
import base64
from datetime import datetime, timezone
import hashlib
import html
import json
import os
from pathlib import Path
import re


ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = ROOT / "outputs" / "current"
SAFE_STEM = re.compile(r"^[A-Za-z0-9][A-Za-z0-9._-]{0,79}$")
SAFE_FORMAT = re.compile(r"^[A-Za-z0-9]{1,12}$")


def files(base: Path) -> list[Path]:
    found: list[Path] = []
    for current, directories, names in os.walk(base, followlinks=False):
        current_path = Path(current)
        directories[:] = [
            name for name in directories if not (current_path / name).is_symlink()
        ]
        for name in names:
            path = current_path / name
            if path.is_file() and not path.is_symlink():
                found.append(path)
    return sorted(found)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def has_format(output_format: str) -> bool:
    aliases = {
        "jpg": {".jpg", ".jpeg"},
        "jpeg": {".jpg", ".jpeg"},
        "tif": {".tif", ".tiff"},
        "tiff": {".tif", ".tiff"},
        "html": {".html", ".htm"},
    }
    suffixes = aliases.get(output_format, {f".{output_format}"})
    return any(path.suffix.lower() in suffixes for path in files(OUTPUT_DIR))


def raster_preview() -> Path:
    candidates = files(OUTPUT_DIR)
    for suffix in (".png", ".jpg", ".jpeg", ".webp"):
        for path in candidates:
            if path.suffix.lower() == suffix:
                return path
    raise RuntimeError("A raster preview is required for HTML and PPTX export")


def validate_raster(path: Path) -> None:
    from PIL import Image

    if path.stat().st_size > 50 * 1024 * 1024:
        raise RuntimeError("The raster preview exceeds the 50 MB preview limit")
    Image.MAX_IMAGE_PIXELS = 100_000_000
    with Image.open(path) as image:
        image.verify()
    with Image.open(path) as image:
        width, height = image.size
    if width <= 0 or height <= 0 or width * height > 100_000_000:
        raise RuntimeError("The raster preview dimensions are unsafe")


def make_html(stem: str, title: str) -> None:
    preview = raster_preview()
    mime = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
    }[preview.suffix.lower()]
    encoded = base64.b64encode(preview.read_bytes()).decode("ascii")
    visual = f'<img src="data:{mime};base64,{encoded}" alt="{html.escape(title)}">'
    target = OUTPUT_DIR / f"{stem}.html"
    target.write_text(
        "<!doctype html>\n"
        '<html lang="en"><head><meta charset="utf-8">'
        f"<title>{html.escape(title)}</title>"
        '<meta name="viewport" content="width=device-width,initial-scale=1">'
        "<style>body{margin:0;background:#f3f5f3;font-family:Arial,sans-serif}"
        "main{max-width:1800px;margin:auto;padding:24px}h1{font-size:20px}"
        "figure{margin:0;background:white;padding:16px}img{display:block;width:100%;height:auto}"
        "</style></head><body><main>"
        f"<h1>{html.escape(title)}</h1><figure>{visual}</figure>"
        "</main></body></html>\n",
        encoding="utf-8",
    )


def make_pptx(stem: str, title: str) -> None:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Pt

    preview = raster_preview()
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title_box = slide.shapes.add_textbox(
        Inches(0.32), Inches(0.16), Inches(12.69), Inches(0.46)
    )
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(18)
    paragraph.font.bold = True
    with Image.open(preview) as image:
        image.verify()
    with Image.open(preview) as image:
        width, height = image.size
    if width <= 0 or height <= 0 or width * height > 100_000_000:
        raise RuntimeError("The preview dimensions are unsafe for PPTX export")
    max_width, max_height = 12.70, 6.62
    scale = min(max_width / width, max_height / height)
    image_width, image_height = width * scale, height * scale
    left = (13.333333 - image_width) / 2
    top = 0.70 + (6.62 - image_height) / 2
    slide.shapes.add_picture(
        str(preview),
        Inches(left),
        Inches(top),
        width=Inches(image_width),
        height=Inches(image_height),
    )
    presentation.save(OUTPUT_DIR / f"{stem}.pptx")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--formats", nargs="+", required=True)
    args = parser.parse_args()
    requested = [value.lower().lstrip(".") for value in args.formats]
    if any(not SAFE_FORMAT.fullmatch(value) for value in requested):
        raise RuntimeError("Invalid requested format")
    project_path = ROOT / "project.json"
    if project_path.stat().st_size > 1024 * 1024:
        raise RuntimeError("The project manifest is too large")
    project = json.loads(project_path.read_text(encoding="utf-8"))
    stem = str(project.get("output_stem", "figure"))
    if not SAFE_STEM.fullmatch(stem):
        raise RuntimeError("The project output stem is unsafe")
    title = str(project.get("title", "Figure Studio output"))[:500]
    raster_preview()
    raster_suffixes = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
    for raster in files(OUTPUT_DIR):
        if raster.suffix.lower() in raster_suffixes:
            validate_raster(raster)
    additions: list[str] = []
    if "html" in requested and not has_format("html"):
        make_html(stem, title)
        additions.append("Created a static HTML export from the raster preview.")
    if "pptx" in requested and not has_format("pptx"):
        make_pptx(stem, title)
        additions.append("Created a one-slide PPTX from the raster preview.")
    missing = [value for value in requested if not has_format(value)]
    if missing:
        raise RuntimeError("Missing requested formats " + ", ".join(missing))
    provenance_path = ROOT / "provenance.json"
    try:
        provenance = json.loads(provenance_path.read_text(encoding="utf-8"))
    except (FileNotFoundError, json.JSONDecodeError):
        provenance = {}
    provenance["generated_at_utc"] = datetime.now(timezone.utc).isoformat()
    provenance["requested_formats"] = requested
    provenance["studio_postprocessing"] = additions
    provenance["outputs"] = [
        {
            "path": str(path.relative_to(ROOT)),
            "bytes": path.stat().st_size,
            "sha256": sha256(path),
        }
        for path in files(OUTPUT_DIR)
    ]
    provenance_path.write_text(
        json.dumps(provenance, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    for addition in additions:
        print(addition)


if __name__ == "__main__":
    main()
